"""
Seed the Learn UI (lesson viewer + practice) from the two workshop decks.

- Renders every slide to a PNG (no LibreOffice needed — drawn with Pillow from
  the slide's own text) and stores it as a LessonSlide.
- Groups slides into Materials (lessons) by SESSION / MODULE markers.
- Pulls a one-line "key point" per slide for the transcript panel.
- Seeds two required drills matching the practice screen.

Run once:  python3 seed_learn_content.py
Idempotent: re-running refreshes slides/keypoints, keeps progress/attempts.
"""
import os, re, json, textwrap
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

os.chdir(os.path.dirname(os.path.abspath(__file__)))
from app import app, db
from models import (Course, Material, LessonSlide, Quiz, QuizQuestion,
                    QuizChoice, User, Enrollment)

DECKS = {
    'vibe': {
        'file': os.path.expanduser('~/Desktop/AI WORKSHOPS/Copy-of-Vibe-Coding-Final.pptx'),
        'title': 'Vibe Coding',
        'programme': 'Vibe Coding Bootcamp',
        'icon': 'fa-terminal',
        'tagline': 'From first prompt to live deployment — AI tools, agents, frontend, and backend.',
        'marker': re.compile(r'^\s*SESSION\s+(\d+)', re.I),
        'sessions': 8,
    },
    'workplace': {
        'file': os.path.expanduser('~/Desktop/AI WORKSHOPS/AI-Workplace (2).pptx'),
        'title': 'AI Workplace',
        'programme': 'AI Workplace Practical Course',
        'icon': 'fa-briefcase',
        'tagline': 'AI marketing copy, decks & video, and building a website with Claude.',
        'marker': re.compile(r'^\s*MODULE\s+0?(\d+)', re.I),
        'sessions': 6,
    },
}

SLIDE_DIR = os.path.join('uploads', 'materials', 'learn')
os.makedirs(SLIDE_DIR, exist_ok=True)

W, H = 1280, 720
FONT = '/System/Library/Fonts/Supplemental/Arial.ttf'
FONT_B = '/System/Library/Fonts/Supplemental/Arial Bold.ttf'
def _f(path, size):
    try: return ImageFont.truetype(path, size)
    except Exception: return ImageFont.load_default()

BG, CARD, INK, SUB, ACCENT = (10, 15, 13), (17, 24, 20), (232, 239, 234), (159, 176, 166), (20, 192, 138)

# Font Awesome 5.15 names (this project ships FA5, not FA6)
ICONS = ['fa-compass', 'fa-lightbulb', 'fa-cube', 'fa-code', 'fa-server',
         'fa-rocket', 'fa-magic', 'fa-project-diagram']


def slide_texts(slide):
    """(title, [body lines]) from a slide's text frames, largest first."""
    frames = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            try: top = Emu(sh.top).inches if sh.top is not None else 99
            except Exception: top = 99
            frames.append((top, sh.text_frame.text.strip()))
    frames.sort(key=lambda x: x[0])
    if not frames:
        return '', []
    title = frames[0][1].splitlines()[0].strip()
    body = []
    for _, t in frames:
        for ln in t.splitlines():
            ln = ln.strip()
            if ln and ln != title and len(ln) > 1:
                body.append(ln)
    return title, body[:7]


def keypoint(title, body):
    base = title
    if body:
        base = f'{title} — {body[0]}' if title else body[0]
    base = re.sub(r'\s+', ' ', base).strip()
    return textwrap.shorten(base, width=180, placeholder='…')


def render_slide(path, kicker, title, body):
    img = Image.new('RGB', (W, H), BG)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle([40, 40, W - 40, H - 40], radius=22, fill=CARD)
    d.rectangle([40, 40, 46, H - 40], fill=ACCENT)
    x = 84
    if kicker:
        d.text((x, 74), kicker.upper(), font=_f(FONT_B, 18), fill=ACCENT)
    fT = _f(FONT_B, 44)
    ty = 116
    for ln in textwrap.wrap(title or 'Slide', width=32)[:3]:
        d.text((x, ty), ln, font=fT, fill=INK); ty += 54
    ty += 16
    fB = _f(FONT, 24)
    for b in body:
        for i, ln in enumerate(textwrap.wrap(b, width=64)[:2]):
            prefix = '•  ' if i == 0 else '   '
            d.text((x, ty), prefix + ln, font=fB, fill=SUB if i else INK)
            ty += 34
        ty += 8
        if ty > H - 90:
            break
    img.save(path, 'PNG')


def build_course(key, cfg):
    if not os.path.exists(cfg['file']):
        print(f'  !! missing deck: {cfg["file"]}')
        return
    prs = Presentation(cfg['file'])
    slides = list(prs.slides)

    course = Course.query.filter_by(title=cfg['title']).first()
    if not course:
        course = Course(title=cfg['title'])
        db.session.add(course)
    course.programme = cfg['programme']
    course.icon = cfg['icon']
    course.tagline = cfg['tagline']
    course.total_sessions = cfg['sessions']
    course.current_session = cfg['sessions']       # everything visible for the demo
    course.description = cfg['tagline']
    db.session.flush()

    # wipe old learn materials for a clean re-seed
    for m in Material.query.filter_by(course_id=course.id).all():
        if (m.filename or '').startswith('learn/'):
            db.session.delete(m)
    db.session.flush()

    # split into lessons: a new lesson starts the FIRST time we see a
    # SESSION n / MODULE n marker (later slides for the same n just append).
    lessons, seen = [], set()
    cur = {'session': 0, 'title': f'{cfg["title"]} — Overview', 'slides': []}
    for s in slides:
        title, body = slide_texts(s)
        mk = cfg['marker'].match(title or '')
        if mk:
            n = int(mk.group(1))
            clean = re.split(r'[·—]', title)[-1].strip()
            clean = re.sub(r'^(SESSION|MODULE)\s+0?\d+\s*[（(].*?[)）]?\s*$', '', clean, flags=re.I).strip()
            clean = re.sub(r'^(SESSION|MODULE)\s+0?\d+\s*', '', clean, flags=re.I).strip()
            if not clean or len(clean) < 4:
                clean = next((re.split(r'[·—]', b)[-1].strip() for b in body
                              if len(b) > 6 and not cfg['marker'].match(b)), '')
            label = 'Session' if key == 'vibe' else 'Module'
            head = f'{label} {n}' + (f' — {clean}' if clean else '')
            if n not in seen:
                seen.add(n)
                if cur['slides']:
                    lessons.append(cur)
                cur = {'session': n, 'title': head, 'slides': []}
        cur['slides'].append((title, body))
    if cur['slides']:
        lessons.append(cur)

    first_mat = None
    for li, lesson in enumerate(lessons):
        sess = max(1, lesson['session'] or 1)
        mat = Material(
            course_id=course.id, session=sess,
            title=lesson['title'][:200],
            description=cfg['tagline'],
            filename=f'learn/{key}-lesson-{li}.png',   # placeholder; real deck view = slides
            file_type='png',
            icon=ICONS[li % len(ICONS)],
            duration_label=f'{len(lesson["slides"])} slides',
            order_index=li,
            is_published=True,
        )
        db.session.add(mat)
        db.session.flush()
        if first_mat is None:
            first_mat = mat

        kps = []
        for si, (stitle, sbody) in enumerate(lesson['slides']):
            png = f'{key}-{li}-{si}.png'
            render_slide(os.path.join(SLIDE_DIR, png),
                         f'{cfg["title"]} · {lesson["title"]}', stitle, sbody)
            kp = keypoint(stitle, sbody)
            kps.append(kp)
            db.session.add(LessonSlide(
                material_id=mat.id, order_index=si,
                image_url=f'/uploads/materials/learn/{png}',
                caption=kp,
            ))
        mat.keypoints = json.dumps(kps)
        # cover image = first slide
        mat.filename = f'learn/{key}-0-0.png' if li == 0 else f'learn/{key}-{li}-0.png'

    db.session.flush()
    print(f'  {cfg["title"]}: course #{course.id}, {len(lessons)} lessons, '
          f'{sum(len(l["slides"]) for l in lessons)} slides')
    return course, first_mat


def seed_drill(course, material, title, desc, q_text, options, correct_idx, explanation):
    quiz = Quiz.query.filter_by(course_id=course.id, title=title).first()
    if not quiz:
        quiz = Quiz(course_id=course.id, title=title)
        db.session.add(quiz)
    quiz.description = desc
    quiz.session = material.session
    quiz.pass_score = 100
    quiz.max_attempts = 5
    quiz.is_published = True
    quiz.is_required = True
    quiz.gates_material_id = material.id
    db.session.flush()
    for qq in list(quiz.questions):
        db.session.delete(qq)
    db.session.flush()
    qq = QuizQuestion(quiz_id=quiz.id, question_text=q_text,
                      question_type='mcq', points=1, explanation=explanation,
                      order_index=0)
    db.session.add(qq)
    db.session.flush()
    for i, opt in enumerate(options):
        db.session.add(QuizChoice(question_id=qq.id, choice_text=opt,
                                  is_correct=(i == correct_idx)))
    print(f'  drill: "{title}" gates material #{material.id}')


with app.app_context():
    db.create_all()
    made = {}
    for key, cfg in DECKS.items():
        print(f'\n{cfg["title"]}  ({cfg["file"].split("/")[-1]})')
        res = build_course(key, cfg)
        if res:
            made[key] = res

    # enrol the demo student in both
    stu = User.query.filter_by(email='student@codencode.my').first()
    if stu:
        for course, _ in made.values():
            if not Enrollment.query.filter_by(student_id=stu.id, course_id=course.id).first():
                db.session.add(Enrollment(student_id=stu.id, course_id=course.id,
                                          payment_status='paid'))
        print(f'\n  enrolled {stu.email} in {len(made)} courses')

    db.session.commit()

    if 'vibe' in made:
        course, first = made['vibe']
        # gate a mid lesson so the drill isn't on the overview
        target = (Material.query.filter_by(course_id=course.id)
                  .filter(Material.session >= 1)
                  .order_by(Material.order_index).all())
        gate_mat = target[1] if len(target) > 1 else first
        seed_drill(
            course, gate_mat,
            'Debug with system prompts',
            "Pick the option that best explains why the model keeps editing files outside the intended scope.",
            "[code]system: \"You are a coding assistant. Help the user improve their project.\"[/code]What's wrong with this system prompt?",
            ["No boundary is set on which files or folders the model may touch",
             "The prompt is too short to be parsed correctly",
             "\"Coding assistant\" is not a recognized role name",
             "System prompts cannot mention the word \"project\""],
            0,
            "A usable system prompt names the folder structure, the naming conventions, and what the model must never touch. Without an explicit boundary the model treats the whole repo as fair game.",
        )

    if 'workplace' in made:
        course, first = made['workplace']
        target = (Material.query.filter_by(course_id=course.id)
                  .filter(Material.session >= 1)
                  .order_by(Material.order_index).all())
        gate_mat = target[0] if target else first
        seed_drill(
            course, gate_mat,
            'Bad / good / perfect prompts',
            "Event: AI for Workplace · 2026-08-22 · 3:00–6:30pm.",
            "Which prompt is closest to \"perfect\" quality for an event poster?",
            ["\"Help me make an event poster.\"",
             "\"Make a poster for the AI for Workplace event on 2026-8-22, 3pm to 6:30pm.\"",
             "\"You are a world-class designer. Make an Instagram-sized poster for AI for Workplace, 2026-8-22, 3–6:30pm, in our mint-teal brand colors, bold modern style, tagline 'work smarter with AI', logo bottom-right.\"",
             "\"Make something nice for an AI event, use whatever colors work.\""],
            2,
            "The perfect prompt gives the model an identity, the exact event facts, brand colors, style direction, the tagline, logo placement, and the output size. Each missing detail is a decision the model makes for you.",
        )

    db.session.commit()
    print('\nDone.')
